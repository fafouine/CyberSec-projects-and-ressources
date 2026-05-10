"""
PowerPoint metadata handler for presentation files.

This module provides the PowerpointHandler class which implements the MetadataHandler
interface for PowerPoint presentation files (.pptx, .pptm, .potx, .potm). Uses python-pptx
for reading and writing presentation properties.

Note:
    Does not support password-protected/encrypted presentations.
"""

import shutil
from pathlib import Path
from typing import Any

from pptx import Presentation

from src.services.metadata_handler import MetadataHandler
from src.utils.exceptions import (
    MetadataNotFoundError,
    MetadataProcessingError,
    MetadataReadingError,
    UnsupportedFormatError,
)

# Supported PowerPoint formats
FORMAT_MAP = {
    "pptx": "pptx",
    "pptm": "pptm",
    "potx": "potx",
    "potm": "potm",
}

# Properties to preserve (not deleted during wipe)
PRESERVED_PROPERTIES = {"created", "modified", "language", "last_printed", "revision"}

# Core properties available in PowerPoint presentations
CORE_PROPERTIES = [
    "author",
    "category",
    "comments",
    "content_status",
    "created",
    "identifier",
    "keywords",
    "language",
    "last_modified_by",
    "last_printed",
    "modified",
    "revision",
    "subject",
    "title",
    "version",
]


class PowerpointHandler(MetadataHandler):
    """
    PowerPoint metadata handler for presentation files.

    Handles extraction and removal of document properties from PowerPoint presentations
    including author, title, subject, keywords, and other core properties.

    Attributes:
        keys_to_delete: List of property names to be wiped.
    """
    def __init__(self, filepath: str):
        """
        Initialize the PowerPoint handler.

        Args:
            filepath: Path to the PowerPoint file to process.
        """
        super().__init__(filepath)
        self.keys_to_delete: list[str] = []

    def _detect_format(self) -> str:
        """
        Detect PowerPoint format from file extension.

        Returns:
            Normalized format string ('pptx', 'pptm', 'potx', or 'potm').

        Raises:
            UnsupportedFormatError: If file extension is not a supported PowerPoint format.
        """
        ext = Path(self.filepath).suffix.lower()
        normalised = FORMAT_MAP.get(ext[1 :])  # Remove leading dot
        if normalised is None:
            raise UnsupportedFormatError(f"Unsupported format: {ext}")

        return normalised

    def read(self) -> dict[str, Any]:
        """
        Extract metadata properties from the PowerPoint presentation.

        Reads all document properties from the presentation and identifies
        which properties should be wiped (excludes created, modified, language, etc.).

        Returns:
            Dictionary of property names to their values.

        Raises:
            MetadataNotFoundError: If no properties are found.
        """
        self.metadata.clear()
        self.keys_to_delete.clear()
        prs = Presentation(str(Path(self.filepath)))
        try:
            if prs.core_properties is None:
                raise MetadataNotFoundError("No metadata found in the file.")

            for attr in CORE_PROPERTIES:
                if hasattr(prs.core_properties, attr):
                    self.metadata[attr] = getattr(prs.core_properties, attr)
                    if attr not in PRESERVED_PROPERTIES:
                        self.keys_to_delete.append(attr)

            return self.metadata
        except Exception as e:
            raise MetadataReadingError(f"error reading metadata. {e}")
        finally:
            del prs

    def wipe(self) -> None:
        """
        Remove metadata properties from the PowerPoint presentation.

        Clears all properties identified during read() except for
        preserved properties (created, modified, language, etc.).

        Raises:
            MetadataNotFoundError: If no properties are found.
        """
        self.processed_metadata.clear()
        prs = Presentation(str(Path(self.filepath)))
        try:
            if prs.core_properties is None:
                raise MetadataNotFoundError("No metadata found in the file.")

            # Clear each property marked for deletion
            for attr in self.keys_to_delete:
                self.processed_metadata[attr] = None
        except Exception as e:
            raise MetadataProcessingError(f"error processing metadata. {e}")
        finally:
            del prs

    def save(self, output_path: str | None = None) -> None:
        """
        Save the presentation with cleaned metadata to the output path.

        Creates a copy of the original file and applies the wiped
        metadata properties to it.

        Args:
            output_path: Path where the cleaned file should be saved.

        Raises:
            ValueError: If output_path is None or empty.
        """
        if not output_path:
            raise ValueError("output_path is required")

        destination_file_path = Path(output_path)
        shutil.copy2(self.filepath, destination_file_path)

        prs = Presentation(str(Path(destination_file_path)))
        try:
            # Apply wiped properties
            for attr in self.processed_metadata:
                if hasattr(prs.core_properties, attr):
                    setattr(prs.core_properties, attr, self.processed_metadata[attr])

            prs.save(str(destination_file_path))
        except Exception as e:
            raise MetadataProcessingError(f"error processing metadata. {e}")
        finally:
            del prs
